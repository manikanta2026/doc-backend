from flask import Flask, request, jsonify
from flask_cors import CORS
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import google.generativeai as genai
import os
from dotenv import load_dotenv  # For loading environment variables

# Load environment variables from .env file
load_dotenv()

# Configure the API key from environment variables
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    raise ValueError("API key for Generative AI is not set. Please ensure that 'GOOGLE_API_KEY' is set in your .env file or as an environment variable.")

genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-2.0-flash")

app = Flask(__name__)
CORS(app)

def extract_text(file):
    filename = file.filename.lower()
    
    if filename.endswith(".pdf"):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    
    elif filename.endswith(".docx"):
        doc = Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    elif filename.endswith(".pptx"):
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    else:
        return None

def generate_summary(text, summary_type):
    if summary_type == "small":
        prompt = (
            f"Summarize the following text in a short bullet-point summary. "
            f"Use '•' for bullet points and ensure each point starts on a new line:\n\n{text}"
        )
    elif summary_type == "medium":
        prompt = (
            f"Provide a medium-length bullet-point summary of the following text. "
            f"If the content is too short, expand with more context. "
            f"Use '•' for bullet points and ensure each point starts on a new line:\n\n{text}"
        )
    else:  # large
        prompt = (
            f"Provide a detailed bullet-point summary of the following text. "
            f"If the content is too short, expand with comprehensive details. "
            f"Use '•' for bullet points and ensure each point starts on a new line:\n\n{text}"
        )
    
    response = model.generate_content(prompt)
    
    if response:
        summary = response.text.strip()
        lines = summary.split('\n')
        formatted_lines = []

        for line in lines:
            line = line.strip()
            if line:
                line = line.lstrip('•*-').strip()
                
                parts = line.split('**')
                formatted_line = ''
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        formatted_line += f'<strong>{part}</strong>'
                    else:
                        formatted_line += part
        
                formatted_line = '• ' + formatted_line
                # Add extra spacing after each bullet point
                formatted_lines.append(formatted_line + '<br><br>')

    formatted_summary = ''.join(formatted_lines)
    return formatted_summary

    return "Summary generation failed."

def generate_qa(text, level="small"):
    if level == "small":
        prompt = (
            f"Generate a few (2-3) high-level, abstract question-answer pairs from the following text strictly in the format:\n\n"
            f"Question:\n[Question here]\nAnswer:\n[Answer here]\n\n"
            f"Ensure each question and answer are on separate lines. Focus on the main ideas only.\n\n{text}"
        )
    elif level == "medium":
        prompt = (
            f"Generate several (4-6) question-answer pairs from the following text strictly in the format:\n\n"
            f"Question:\n[Question here]\nAnswer:\n[Answer here]\n\n"
            f"Ensure each question and answer are on separate lines. Include both abstract and some detailed questions.\n\n{text}"
        )
    else:  # large
        prompt = (
            f"Generate many (7 or more) detailed question-answer pairs from the following text strictly in the format:\n\n"
            f"Question:\n[Question here]\nAnswer:\n[Answer here]\n\n"
            f"Ensure each question and answer are on separate lines. Cover all levels: abstract, intermediate, and detailed questions.\n\n{text}"
        )
    response = model.generate_content(prompt)
    if response:
        text = response.text.strip()
        lines = text.split("\n")
        formatted_text = []

        current_block = []
        for line in lines:
            line = line.strip()
            if line.startswith("Question:"):
                # If there's a previous block, add it with spacing
                if current_block:
                    formatted_text.append("".join(current_block) + "<br><br>")
                    current_block = []
                current_block.append(f"<br><strong>{line}</strong>")
            elif line.startswith("Answer:"):
                current_block.append(f"<br><strong>{line}</strong>")
            else:
                current_block.append(line)
        # Add the last block if present
        if current_block:
            formatted_text.append("".join(current_block) + "<br><br>")

        formatted_response = "".join(formatted_text).strip()
        return formatted_response
    else:
        return "Response generation failed."

@app.route('/summary', methods=['POST'])
def summarize_file():
    file = request.files.get('file')
    summary_type = request.form.get('summary_type', 'small')
    
    if not file:
        return jsonify({"error": "No file provided"}), 400

    text = extract_text(file)
    if not text:
        return jsonify({"error": "Failed to extract text from the uploaded file. Make sure it’s a .pdf, .docx, or .pptx"}), 500

    result = generate_summary(text, summary_type)
    return jsonify({"summary": result})

@app.route('/qa', methods=['POST'])
def qa_file():
    file = request.files.get('file')
    level = request.form.get('summary_type', 'small')
    
    if not file:
        return jsonify({"error": "No file provided"}), 400

    text = extract_text(file)
    if not text:
        return jsonify({"error": "Failed to extract text from the uploaded file. Make sure it’s a .pdf, .docx, or .pptx"}), 500

    result = generate_qa(text, level)
    return jsonify({"qa": result})


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=True)
